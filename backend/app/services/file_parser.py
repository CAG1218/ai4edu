"""
AI4Edu 文件解析器
支持 PDF / DOCX / PPTX 等格式的文本提取
"""
import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)


class FileParser:
    """文件解析器"""

    @staticmethod
    async def parse_pdf(content: bytes) -> str:
        """
        解析 PDF 文件提取文本

        Args:
            content: 文件二进制内容

        Returns:
            提取的文本内容
        """
        try:
            import pdfplumber

            text_parts = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            return "\n".join(text_parts)
        except ImportError:
            logger.warning("pdfplumber 未安装，尝试 PyPDF2")
            try:
                from PyPDF2 import PdfReader

                reader = PdfReader(io.BytesIO(content))
                text_parts = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                return "\n".join(text_parts)
            except ImportError:
                logger.error("PDF 解析库未安装")
                return ""
        except Exception as e:
            logger.error(f"PDF 解析失败: {e}")
            return ""

    @staticmethod
    async def parse_docx(content: bytes) -> str:
        """
        解析 DOCX 文件提取文本

        Args:
            content: 文件二进制内容

        Returns:
            提取的文本内容
        """
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n".join(paragraphs)
        except ImportError:
            logger.error("python-docx 未安装")
            return ""
        except Exception as e:
            logger.error(f"DOCX 解析失败: {e}")
            return ""

    @staticmethod
    async def parse_pptx(content: bytes) -> str:
        """
        解析 PPTX 文件提取文本

        Args:
            content: 文件二进制内容

        Returns:
            提取的文本内容
        """
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except ImportError:
            logger.error("python-pptx 未安装")
            return ""
        except Exception as e:
            logger.error(f"PPTX 解析失败: {e}")
            return ""

    @staticmethod
    async def parse_file(content: bytes, mime_type: Optional[str] = None, filename: Optional[str] = None) -> str:
        """
        自动路由：根据文件类型选择合适的解析器

        Args:
            content: 文件二进制内容
            mime_type: MIME 类型
            filename: 文件名

        Returns:
            提取的文本内容
        """
        # 确定文件类型
        ext = ""
        if filename:
            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        elif mime_type:
            mime_map = {
                "application/pdf": "pdf",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
            }
            ext = mime_map.get(mime_type, "")

        parser_map = {
            "pdf": FileParser.parse_pdf,
            "docx": FileParser.parse_docx,
            "pptx": FileParser.parse_pptx,
        }

        parser = parser_map.get(ext)
        if parser:
            return await parser(content)

        # 对于文本类文件直接解码
        text_exts = {"txt", "md", "csv", "json", "xml", "html", "css", "js", "py", "java"}
        if ext in text_exts:
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                try:
                    return content.decode("gbk")
                except UnicodeDecodeError:
                    return ""

        logger.warning(f"不支持的文件类型: {ext or mime_type}")
        return ""
